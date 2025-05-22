import os
import subprocess
import logging
import tempfile
from typing import Optional
import time

# 设置日志
logger = logging.getLogger(__name__)

# 尝试导入所有转换所需的库
try:
    from pdf2docx import Converter as PDFToDocxConverter
except ImportError:
    logger.warning("无法导入pdf2docx库，PDF转Word功能可能不可用")

try:
    import tabula
    import pandas as pd
except ImportError:
    logger.warning("无法导入tabula或pandas库，PDF转Excel功能可能不可用")

try:
    from pdf2image import convert_from_path
except ImportError:
    logger.warning("无法导入pdf2image库，PDF转图片功能可能不可用")

try:
    from PIL import Image
except ImportError:
    logger.warning("无法导入PIL库，图片处理功能可能不可用")

try:
    import PyPDF2
except ImportError:
    logger.warning("无法导入PyPDF2库，PDF功能可能受限")

try:
    from docx import Document
except ImportError:
    logger.warning("无法导入python-docx库，Word文档处理功能可能不可用")

try:
    from fpdf import FPDF
except ImportError:
    logger.warning("无法导入fpdf库，文本转PDF功能可能不可用")

try:
    import markdown
except ImportError:
    logger.warning("无法导入markdown或weasyprint库，Markdown转PDF功能可能不可用")


def convert_file(input_path: str, to_format: str, output_path: str, quality: int = 2, original_filename: str = None) -> dict:
    """
    转换文件从一种格式到另一种格式

    Args:
        input_path: 输入文件路径
        to_format: 目标格式
        output_path: 输出文件路径
        quality: 质量等级 (1=低, 2=中, 3=高)
        original_filename: 原始上传的文件名

    Returns:
        包含转换结果和原始文件名信息的字典
    """
    input_ext = os.path.splitext(input_path)[1][1:].lower()

    # 如果没有提供原始文件名，使用输入路径的文件名
    if original_filename is None:
        original_filename = os.path.basename(input_path)
    
    # 记录原始文件名
    logger.info(f"开始转换: {original_filename} ({input_ext} -> {to_format}), 质量级别: {quality}")

    # 最终转换结果
    result = {
        "output_path": "",
        "original_filename": original_filename,
        "output_filename": os.path.basename(output_path),
        "input_format": input_ext,
        "output_format": to_format
    }

    # 转换函数返回值处理
    def process_result(res):
        """处理转换函数的返回值，确保返回字符串路径"""
        if res is None:
            logger.error("转换函数返回了None")
            raise ValueError("转换结果为空")
            
        if isinstance(res, dict):
            if "output_path" in res:
                return res["output_path"]
            else:
                logger.error(f"转换函数返回了不包含output_path的字典: {res}")
                raise ValueError("转换结果字典缺少output_path字段")
                
        if not isinstance(res, (str, bytes, os.PathLike)):
            logger.error(f"转换函数返回了不受支持的类型: {type(res)}")
            raise ValueError(f"转换结果类型不支持: {type(res)}")
            
        return res

    # PDF转其他格式
    if input_ext == 'pdf':
        if to_format == 'docx':
            result["output_path"] = process_result(pdf_to_docx(input_path, output_path, quality))
        elif to_format == 'xlsx':
            result["output_path"] = process_result(pdf_to_excel(input_path, output_path, quality))
        elif to_format == 'pptx':
            result["output_path"] = process_result(pdf_to_pptx(input_path, output_path, quality))
        elif to_format in ['jpg', 'png']:
            result["output_path"] = process_result(pdf_to_images(input_path, output_path, quality, to_format))
        elif to_format in ['dwg', 'dxf', 'cad']:
            result["output_path"] = process_result(pdf_to_cad(input_path, output_path, quality))
        elif to_format in ['scannable_pdf', 'scanned_pdf']:
            # 使用提供的输出路径，不做修改
            result["output_path"] = process_result(pdf_to_scannable_pdf(input_path, output_path, quality))
            # 更新结果中的输出文件名
            result["output_filename"] = os.path.basename(output_path)
        elif to_format in ['searchable_pdf']:
            # 使用提供的输出路径，不做修改
            result["output_path"] = process_result(pdf_to_searchable_pdf(input_path, output_path, quality))
            # 更新结果中的输出文件名
            result["output_filename"] = os.path.basename(output_path)

    # 其他格式转PDF
    elif to_format == 'pdf':
        if input_ext in ['jpg', 'jpeg', 'png']:
            result["output_path"] = process_result(image_to_pdf(input_path, output_path, quality))
        elif input_ext in ['doc', 'docx']:
            result["output_path"] = process_result(word_to_pdf(input_path, output_path, quality))
        elif input_ext in ['xls', 'xlsx']:
            result["output_path"] = process_result(excel_to_pdf(input_path, output_path, quality))
        elif input_ext in ['ppt', 'pptx']:
            result["output_path"] = process_result(ppt_to_pdf(input_path, output_path, quality))
        elif input_ext == 'txt':
            result["output_path"] = process_result(txt_to_pdf(input_path, output_path, quality))
        elif input_ext == 'md':
            result["output_path"] = process_result(markdown_to_pdf(input_path, output_path, quality))

    if not result["output_path"]:
        raise ValueError(f"不支持从 {input_ext} 转换到 {to_format}")
    
    # 最终确保输出路径是一个有效的字符串路径，而不是字典
    if isinstance(result["output_path"], dict):
        if "output_path" in result["output_path"]:
            result["output_path"] = result["output_path"]["output_path"]
        else:
            logger.error(f"转换结果包含无效的嵌套字典: {result['output_path']}")
            raise ValueError("转换结果包含无效的嵌套字典")
    
    # 验证输出路径的存在性和类型
    if not isinstance(result["output_path"], (str, bytes, os.PathLike)):
        logger.error(f"转换函数返回了无效的输出路径类型: {type(result['output_path'])}")
        raise ValueError(f"转换函数返回了无效的输出路径类型: {type(result['output_path'])}")
    
    # 检查文件是否存在
    if not os.path.exists(result["output_path"]):
        logger.error(f"转换后的文件不存在: {result['output_path']}")
        raise ValueError(f"转换后的文件不存在: {result['output_path']}")
        
    logger.info(f"转换成功完成: {result['output_path']}")
    
    return result


# PDF 转 Word
def pdf_to_docx(input_path: str, output_path: str, quality: int) -> str:
    """将PDF转换为DOCX格式"""
    try:
        logger.info(f"开始PDF转Word转换: {input_path}")
        
        from pdf2docx.converter import Converter
        
        # 根据质量参数设置转换选项
        conversion_options = {
            "debug": False,
            "keep_image_as_whole": True,  # 保持图像完整性
            "layout_analysis_mode": min(2, quality),  # 布局分析强度
            "implicit_section_break": False,  # 避免添加额外分节符
            "fast_layout_mode": quality < 2,  # 低质量时使用快速布局
            "minimize_tables": quality < 3,   # 高质量时保留复杂表格
        }
        
        # 使用高级转换器
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None, **conversion_options)
        cv.close()
        
        logger.info(f"PDF转Word完成: {output_path}")
        
        # 确保处理所有可能的返回路径
        if isinstance(output_path, dict) and "output_path" in output_path:
            return output_path["output_path"]
        return output_path
    except Exception as e:
        logger.error(f"PDF转Word失败: {str(e)}")
        raise ValueError(f"PDF转Word失败: {str(e)}")


# PDF 转 Excel
def pdf_to_excel(input_path: str, output_path: str, quality: int) -> str:
    """将PDF转换为Excel格式"""
    try:
        logger.info(f"开始PDF转Excel转换: {input_path}")
        # 调整表格提取精度
        lattice = quality >= 2  # 如果质量大于等于2，使用lattice模式识别表格线

        # 从PDF中提取表格
        dfs = tabula.read_pdf(
            input_path,
            pages='all',
            lattice=lattice,
            multiple_tables=True
        )

        # 如果没有提取到表格，创建一个空表格
        if not dfs:
            pd.DataFrame().to_excel(output_path, index=False)
            logger.warning(f"PDF中未检测到表格，生成空Excel文件")
            return output_path

        # 创建Excel写入器
        with pd.ExcelWriter(output_path) as writer:
            for i, df in enumerate(dfs):
                sheet_name = f"Sheet{i + 1}"
                df.to_excel(writer, sheet_name=sheet_name, index=False)

        logger.info(f"PDF转Excel完成: {output_path}, 共{len(dfs)}个表格")
        
        # 确保处理所有可能的返回路径
        if isinstance(output_path, dict) and "output_path" in output_path:
            return output_path["output_path"]
        return output_path
    except Exception as e:
        logger.error(f"PDF转Excel失败: {str(e)}")
        raise ValueError(f"PDF转Excel失败: {str(e)}")


# PDF 转 PPT
def pdf_to_pptx(input_path: str, output_path: str, quality: int) -> str:
    """将PDF转换为PPTX格式"""
    try:
        logger.info(f"开始PDF转PPT转换: {input_path}")
        from pptx import Presentation
        from pptx.util import Inches

        # 先将PDF转换为图片
        images = convert_from_path(
            input_path,
            dpi=100 * quality  # 根据质量调整DPI
        )

        logger.info(f"已将PDF转换为{len(images)}张图片")

        # 创建PPT
        prs = Presentation()

        # 为每一页PDF创建一张幻灯片
        for i, img in enumerate(images):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

            # 保存临时图片
            temp_img_path = tempfile.mktemp(suffix='.png')
            img.save(temp_img_path, format='PNG')

            # 添加图片到幻灯片
            slide.shapes.add_picture(temp_img_path, 0, 0, prs.slide_width, prs.slide_height)

            # 删除临时图片
            os.remove(temp_img_path)

            logger.info(f"已处理PDF页面 {i + 1}/{len(images)}")

        prs.save(output_path)
        logger.info(f"PDF转PPT完成: {output_path}")
        
        # 确保处理所有可能的返回路径
        if isinstance(output_path, dict) and "output_path" in output_path:
            return output_path["output_path"]
        return output_path
    except Exception as e:
        logger.error(f"PDF转PPT失败: {str(e)}")
        raise ValueError(f"PDF转PPT失败: {str(e)}")


# PDF 转图片
def pdf_to_images(input_path: str, output_path: str, quality: int, img_format: str) -> str:
    """将PDF转换为图片格式"""
    try:
        logger.info(f"开始PDF转图片转换: {input_path}")
        # 调整DPI基于质量
        dpi = 100 * quality

        # 将PDF转换为图片
        images = convert_from_path(input_path, dpi=dpi)
        logger.info(f"已将PDF转换为{len(images)}张图片")

        # 如果有多页，创建ZIP文件
        if len(images) > 1:
            import zipfile
            from io import BytesIO

            # 更改输出路径为ZIP
            zip_path = os.path.splitext(output_path)[0] + '.zip'

            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for i, img in enumerate(images):
                    img_buffer = BytesIO()
                    img.save(img_buffer, format=img_format.upper())
                    img_buffer.seek(0)
                    zipf.writestr(f"page_{i + 1}.{img_format}", img_buffer.read())

            logger.info(f"已创建包含{len(images)}张图片的ZIP文件: {zip_path}")
            return zip_path
        else:
            # 如果只有一页，直接保存
            images[0].save(output_path, format=img_format.upper())
            logger.info(f"PDF转单张图片完成: {output_path}")
            return output_path
        
        # 确保处理所有可能的返回路径
        if isinstance(output_path, dict) and "output_path" in output_path:
            return output_path["output_path"]
        return output_path
    except Exception as e:
        logger.error(f"PDF转图片失败: {str(e)}")
        raise ValueError(f"PDF转图片失败: {str(e)}")


# PDF 转 CAD
def pdf_to_cad(input_path: str, output_path: str, quality: int) -> str:
    """将PDF转换为CAD格式（DWG/DXF）"""
    try:
        logger.info(f"开始PDF转CAD转换: {input_path}")
        
        # CAD转换通常需要专业软件，这里实现一个提示功能
        # 实际项目中可能需要集成第三方服务或商业软件
        
        # 创建一个简单的文本文件，说明需要专业软件
        output_format = os.path.splitext(output_path)[1][1:].lower()
        text_output_path = os.path.splitext(output_path)[0] + ".txt"
        
        with open(text_output_path, 'w', encoding='utf-8') as f:
            f.write(f"""
PDF转{output_format.upper()}转换需要专业CAD软件支持。

建议使用以下方法：
1. 使用AutoCAD软件打开PDF文件
2. 使用PDF2CAD等专业转换软件
3. 咨询专业CAD服务提供商

文件: {os.path.basename(input_path)}
请求时间: {time.strftime('%Y-%m-%d %H:%M:%S')}
            """)
            
        # 复制原始PDF到输出目录，以便用户可以下载
        import shutil
        pdf_output_path = os.path.splitext(output_path)[0] + ".pdf"
        shutil.copy(input_path, pdf_output_path)
        
        logger.info(f"创建了CAD转换说明文件: {text_output_path}")
        
        # 返回文本文件路径
        return text_output_path
    except Exception as e:
        logger.error(f"PDF转CAD失败: {str(e)}")
        raise ValueError(f"PDF转CAD失败: {str(e)}")


# PDF 转扫描版PDF
def pdf_to_scannable_pdf(input_path: str, output_path: str, quality: int) -> str:
    """将普通PDF转换为看起来像扫描版的PDF"""
    try:
        logger.info(f"开始转换PDF为扫描版: {input_path}")
        logger.info(f"输出路径设置为: {output_path}")
        
        # 将PDF转换为图片
        dpi = 150 * quality  # 调整分辨率
        images = convert_from_path(
            input_path, 
            dpi=dpi,
            grayscale=quality < 3,  # 质量小于3时使用灰度
            transparent=False
        )
        
        logger.info(f"已将PDF转换为{len(images)}张图片")
        
        # 处理图片使其看起来像扫描版
        processed_images = []
        
        for i, img in enumerate(images):
            # 调整图片让它看起来像扫描版
            from PIL import Image, ImageFilter, ImageEnhance
            
            logger.info(f"处理第 {i+1} 张图片，模式: {img.mode}, 尺寸: {img.size}")
            
            # 降低质量
            if quality < 3:
                # 添加模糊
                img = img.filter(ImageFilter.GaussianBlur(radius=0.5))
                
                # 调整对比度
                enhancer = ImageEnhance.Contrast(img)
                img = enhancer.enhance(1.2)
                
                # 添加噪点 - 根据图像模式使用合适的颜色格式
                import random
                logger.info(f"为图片添加噪点，图像模式: {img.mode}")
                for _ in range(3000):
                    x = random.randint(0, img.width - 1)
                    y = random.randint(0, img.height - 1)
                    # 根据图像模式设置不同的颜色值
                    if img.mode == 'L':  # 灰度模式
                        img.putpixel((x, y), 0)  # 0 表示黑色
                    elif img.mode == '1':  # 二值模式
                        img.putpixel((x, y), 0)  # 0 表示黑色
                    else:  # RGB 或其他彩色模式
                        img.putpixel((x, y), (0, 0, 0))  # 黑色
            
            processed_images.append(img)
            logger.info(f"已处理扫描效果 {i + 1}/{len(images)}")
        
        # 合并回PDF - 直接使用提供的输出路径
        logger.info(f"开始生成PDF文件到: {output_path}")
        processed_images[0].save(
            output_path,
            "PDF",
            resolution=100.0,
            save_all=True,
            append_images=processed_images[1:] if len(processed_images) > 1 else []
        )
        
        # 检查生成的文件
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            logger.info(f"生成扫描版PDF完成: {output_path}, 文件大小: {file_size} 字节")
            return output_path
        else:
            raise ValueError(f"无法生成PDF文件: {output_path}")
    except Exception as e:
        logger.error(f"PDF转扫描版失败: {str(e)}")
        raise ValueError(f"PDF转扫描版失败: {str(e)}")


# PDF 转可搜索PDF
def pdf_to_searchable_pdf(input_path: str, output_path: str, quality: int) -> str:
    """将PDF转换为可搜索的PDF（添加OCR文本层）"""
    try:
        logger.info(f"开始转换PDF为可搜索版: {input_path}")
        logger.info(f"输出路径设置为: {output_path}")
        
        # 需要安装pytesseract和tesseract-ocr
        try:
            import pytesseract
            from pdf2image import convert_from_path
            import PyPDF2
            from PIL import Image, ImageDraw, ImageFont
            import io
            
            # OCR库都可用，执行正常的OCR流程
            
            # 检查tesseract是否可用
            try:
                tesseract_version = pytesseract.get_tesseract_version()
                logger.info(f"Tesseract版本: {tesseract_version}")
            except Exception as e:
                logger.error(f"Tesseract不可用: {str(e)}")
                raise ValueError("未安装Tesseract OCR或无法访问。请安装Tesseract并确保它在PATH中。")
            
            # 检查PDF是否已经是可搜索的
            try:
                pdf_reader = PyPDF2.PdfReader(input_path)
                text = pdf_reader.pages[0].extract_text().strip()
                if text:
                    logger.info("PDF已经包含文本层，可能已经是可搜索的")
                    import shutil
                    shutil.copy(input_path, output_path)
                    return output_path
            except Exception as e:
                logger.warning(f"检查PDF文本层时出错: {str(e)}")
            
            # 将PDF转换为图片
            dpi = 300  # 更高的DPI有助于提高OCR精度
            images = convert_from_path(input_path, dpi=dpi)
            logger.info(f"已将PDF转换为{len(images)}张图片")
            
            # 创建一个临时目录用于存储中间文件
            import tempfile
            with tempfile.TemporaryDirectory() as tmpdirname:
                logger.info(f"创建临时目录: {tmpdirname}")
                
                # 准备字体
                font_path = None
                font_name = None
                for font_test in [
                    ("黑体", r"C:\Windows\Fonts\simhei.ttf"),
                    ("宋体", r"C:\Windows\Fonts\simsun.ttc"),
                    ("微软雅黑", r"C:\Windows\Fonts\msyh.ttc")
                ]:
                    name, path = font_test
                    if os.path.exists(path):
                        font_path = path
                        font_name = name
                        logger.info(f"找到中文字体: {name} ({path})")
                        break
                
                if not font_path:
                    logger.warning("未找到中文字体，OCR结果可能无法正确显示")
                
                # 创建一个PDF写入器用于最终PDF
                pdf_writer = PyPDF2.PdfWriter()
                
                for i, img in enumerate(images):
                    logger.info(f"处理第{i+1}页...")
                    
                    # 使用OCR获取文本 - 同时使用中文简体和英文识别
                    text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                    
                    # 预览OCR结果的前100个字符
                    if text and len(text) > 0:
                        preview = text[:min(100, len(text))].replace('\n', ' ')
                        logger.info(f"OCR文本预览: {preview}...")
                        logger.info(f"OCR文本长度: {len(text)}")
                    else:
                        logger.warning("OCR未提取到文本")
                        text = ""
                    
                    # 保存图像为临时图像PDF
                    img_pdf_path = os.path.join(tmpdirname, f"img_{i}.pdf")
                    img.convert('RGB').save(img_pdf_path, "PDF")
                    
                    # 方法1: 使用reportlab创建文本层
                    text_pdf_path = os.path.join(tmpdirname, f"text_{i}.pdf")
                    
                    try:
                        # 改进后的文本层创建方法
                        from reportlab.pdfgen import canvas
                        from reportlab.pdfbase import pdfmetrics
                        from reportlab.pdfbase.ttfonts import TTFont
                        
                        # 注册中文字体
                        font_registered = False
                        if font_path:
                            try:
                                if "simhei" in font_path.lower():
                                    pdfmetrics.registerFont(TTFont('SimHei', font_path))
                                    font_registered = True
                                    font_name = 'SimHei'
                                elif "simsun" in font_path.lower():
                                    pdfmetrics.registerFont(TTFont('SimSun', font_path))
                                    font_registered = True
                                    font_name = 'SimSun'
                                elif "msyh" in font_path.lower():
                                    pdfmetrics.registerFont(TTFont('MicrosoftYaHei', font_path))
                                    font_registered = True
                                    font_name = 'MicrosoftYaHei'
                                logger.info(f"已注册字体: {font_name}")
                            except Exception as font_e:
                                logger.warning(f"注册字体失败: {str(font_e)}")
                        
                        # 获取图像尺寸
                        img_width, img_height = img.size
                        
                        # 创建一个新的canvas用于文本
                        c = canvas.Canvas(text_pdf_path, pagesize=(img_width, img_height))
                        
                        if font_registered:
                            c.setFont(font_name, 10)  # 使用适当大小的字体
                        else:
                            c.setFont("Helvetica", 10)
                        
                        # 使用低不透明度文本 (几乎不可见但可搜索)
                        c.setFillColorRGB(0, 0, 0, 0.01)  
                        
                        # 添加OCR文本
                        y_pos = img_height - 20
                        for line in text.split('\n'):
                            if line.strip():
                                try:
                                    c.drawString(10, y_pos, line)
                                except:
                                    # 如果绘制文本失败，尝试仅使用ASCII字符
                                    ascii_line = ''.join(c if ord(c) < 128 else '_' for c in line)
                                    c.drawString(10, y_pos, ascii_line)
                                    logger.warning(f"使用ASCII模式绘制文本行")
                                y_pos -= 12
                                
                                # 防止文本超出页面底部
                                if y_pos < 10:
                                    y_pos = img_height - 20
                        
                        c.save()
                        logger.info(f"成功创建文本层PDF")
                        
                        # 合并图像PDF和文本层PDF
                        img_pdf = PyPDF2.PdfReader(img_pdf_path)
                        text_pdf = PyPDF2.PdfReader(text_pdf_path)
                        
                        page = img_pdf.pages[0]
                        if len(text_pdf.pages) > 0:
                            page.merge_page(text_pdf.pages[0])
                        
                        pdf_writer.add_page(page)
                        logger.info(f"已合并第{i+1}页的图像和文本层")
                        
                    except Exception as text_err:
                        logger.error(f"创建文本层时出错: {str(text_err)}")
                        
                        # 方法2: 如果reportlab方法失败，尝试PIL+PDF方法
                        try:
                            logger.info("尝试使用PIL图像方法创建文本层")
                            
                            # 创建透明图像作为文本层
                            text_img = Image.new('RGBA', img.size, (255, 255, 255, 0))
                            draw = ImageDraw.Draw(text_img)
                            
                            # 尝试加载中文字体
                            try:
                                if font_path:
                                    font = ImageFont.truetype(font_path, 12)
                                else:
                                    # 尝试默认字体
                                    font = ImageFont.load_default()
                            except Exception as font_err:
                                logger.warning(f"加载字体失败: {str(font_err)}")
                                font = None
                            
                            # 绘制文本到透明图像
                            y_pos = 20
                            for line in text.split('\n'):
                                if line.strip():
                                    try:
                                        # 使用非常淡的颜色 (几乎透明)
                                        if font:
                                            draw.text((10, y_pos), line, fill=(0, 0, 0, 10), font=font)
                                        else:
                                            draw.text((10, y_pos), line, fill=(0, 0, 0, 10))
                                    except Exception as draw_err:
                                        logger.warning(f"绘制文本失败: {str(draw_err)}")
                                    y_pos += 15
                            
                            # 将文本图像转换为PDF
                            text_img_path = os.path.join(tmpdirname, f"text_img_{i}.png")
                            text_img.save(text_img_path, "PNG")
                            
                            # 创建一个包含两个图层的PDF
                            from reportlab.pdfgen import canvas
                            from reportlab.lib.utils import ImageReader
                            
                            combined_pdf_path = os.path.join(tmpdirname, f"combined_{i}.pdf")
                            c = canvas.Canvas(combined_pdf_path, pagesize=img.size)
                            
                            # 先添加原始图像
                            img_path = os.path.join(tmpdirname, f"img_{i}.png")
                            img.save(img_path, "PNG")
                            c.drawImage(ImageReader(img_path), 0, 0, width=img.width, height=img.height)
                            
                            # 在上面添加文本图层
                            c.drawImage(ImageReader(text_img_path), 0, 0, width=img.width, height=img.height)
                            c.save()
                            
                            # 将生成的PDF添加到最终PDF
                            combined_pdf = PyPDF2.PdfReader(combined_pdf_path)
                            pdf_writer.add_page(combined_pdf.pages[0])
                            logger.info(f"已使用图像方法合并第{i+1}页")
                            
                        except Exception as pil_err:
                            logger.error(f"图像方法也失败: {str(pil_err)}")
                            # 如果两种方法都失败，只添加图像层
                            img_pdf = PyPDF2.PdfReader(img_pdf_path)
                            pdf_writer.add_page(img_pdf.pages[0])
                            logger.info(f"只添加了第{i+1}页的图像层（无文本层）")
                
                # 保存最终的PDF
                with open(output_path, 'wb') as f:
                    pdf_writer.write(f)
                
                logger.info(f"已创建可搜索PDF: {output_path}")
                return output_path
            
        except ImportError as e:
            logger.warning(f"缺少OCR库: {str(e)}")
            
            # 创建一个替代解决方案：通知用户并返回原始PDF
            try:
                import shutil
                import PyPDF2
                
                # 创建一个说明PDF
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                
                instructions_path = f"{output_path}.instructions.pdf"
                c = canvas.Canvas(instructions_path, pagesize=letter)
                
                c.setFont("Helvetica-Bold", 16)
                c.drawString(72, 700, "PDF OCR Processing Failed")
                
                c.setFont("Helvetica", 12)
                c.drawString(72, 670, "The PDF could not be made searchable because:")
                c.drawString(72, 650, "1. Required OCR libraries are missing")
                c.drawString(72, 630, "2. The original PDF has been included after this page")
                
                c.drawString(72, 590, "To enable OCR functionality, install:")
                c.drawString(72, 570, "1. Python libraries: pip install pytesseract pdf2image")
                c.drawString(72, 550, "2. Tesseract OCR: https://github.com/tesseract-ocr/tesseract")
                
                # 添加中文说明作为图片
                try:
                    img = Image.new('RGB', (500, 200), color=(255, 255, 255))
                    d = ImageDraw.Draw(img)
                    
                    # 查找中文字体
                    font = None
                    for font_path in [
                        r"C:\Windows\Fonts\simhei.ttf",
                        r"C:\Windows\Fonts\simsun.ttc",
                        r"C:\Windows\Fonts\msyh.ttc"
                    ]:
                        if os.path.exists(font_path):
                            try:
                                font = ImageFont.truetype(font_path, 14)
                                logger.info(f"使用字体: {font_path}")
                                break
                            except Exception as font_err:
                                logger.warning(f"加载字体失败: {str(font_err)}")
                    
                    # 使用找到的字体或默认字体绘制文本
                    if font:
                        d.text((10, 10), "PDF未进行OCR处理", fill=(0, 0, 0), font=font)
                        d.text((10, 40), "由于系统缺少必要的OCR组件，无法创建可搜索的PDF。", fill=(0, 0, 0), font=font)
                        d.text((10, 70), "已返回原始PDF文件，但未添加文本层。", fill=(0, 0, 0), font=font)
                        d.text((10, 100), "如需OCR功能，请安装pytesseract和Tesseract OCR。", fill=(0, 0, 0), font=font)
                    else:
                        d.text((10, 10), "PDF OCR Processing Failed (Chinese support unavailable)", fill=(0, 0, 0))
                    
                    # 将图像添加到PDF
                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format='PNG')
                    img_buffer.seek(0)
                    
                    c.drawImage(img_buffer, 72, 400, width=450, height=130)
                    logger.info("已添加中文说明图像")
                
                except Exception as img_err:
                    logger.warning(f"创建中文说明图像失败: {str(img_err)}")
                
                c.save()
                
                # 复制原始PDF并合并
                temp_original = output_path + ".original.pdf"
                shutil.copy(input_path, temp_original)
                
                # 合并PDF
                merger = PyPDF2.PdfMerger()
                merger.append(instructions_path)
                merger.append(temp_original)
                merger.write(output_path)
                merger.close()
                
                # 清理临时文件
                os.remove(instructions_path)
                os.remove(temp_original)
                
                logger.info(f"已创建包含说明页的PDF(无OCR): {output_path}")
                return output_path
                
            except Exception as merge_err:
                logger.error(f"创建说明页失败: {str(merge_err)}")
                # 如果无法创建说明页，直接返回原始PDF的副本
                import shutil
                shutil.copy(input_path, output_path)
                logger.info(f"已复制原始PDF(无OCR): {output_path}")
                return output_path
                
    except Exception as e:
        logger.error(f"PDF转可搜索版失败: {str(e)}")
        # 尝试返回原始PDF
        try:
            import shutil
            shutil.copy(input_path, output_path)
            logger.info(f"已复制原始PDF作为备选方案: {output_path}")
            return output_path
        except Exception as copy_err:
            logger.error(f"复制原始PDF也失败: {str(copy_err)}")
            raise ValueError(f"PDF转可搜索版失败且无法返回原始PDF: {str(e)}")


# 图片转PDF
def image_to_pdf(input_path: str, output_path: str, quality: int) -> str:
    """将图片转换为PDF格式"""
    try:
        logger.info(f"开始图片转PDF转换: {input_path}")
        image = Image.open(input_path)

        # 根据质量调整图像大小
        if quality < 3:
            width, height = image.size
            new_width = int(width * 0.8 * quality)
            new_height = int(height * 0.8 * quality)
            image = image.resize((new_width, new_height))
            logger.info(f"已根据质量调整图片大小: {width}x{height} -> {new_width}x{new_height}")

        # 如果是单通道图像(如灰度图)，转换为RGB
        if image.mode != 'RGB':
            image = image.convert('RGB')
            logger.info("已将图片转换为RGB模式")

        image.save(output_path, 'PDF', resolution=100.0 * quality)
        logger.info(f"图片转PDF完成: {output_path}")
        
        # 确保处理所有可能的返回路径
        if isinstance(output_path, dict) and "output_path" in output_path:
            return output_path["output_path"]
        return output_path
    except Exception as e:
        logger.error(f"图片转PDF失败: {str(e)}")
        raise ValueError(f"图片转PDF失败: {str(e)}")


# Word 转 PDF
def word_to_pdf(input_path: str, output_path: str, quality: int) -> str:
    """将Word转换为PDF格式"""
    try:
        logger.info(f"开始Word转PDF转换: {input_path}")
        conversion_success = False
        
        # 在Windows上，优先使用Office COM组件
        if os.name == 'nt':
            try:
                import win32com.client
                import pythoncom
                
                logger.info("尝试使用Word COM组件转换")
                pythoncom.CoInitialize()
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                
                # 使用绝对路径
                abs_input_path = os.path.abspath(input_path)
                abs_output_path = os.path.abspath(output_path)
                
                # 确保输出目录存在
                os.makedirs(os.path.dirname(abs_output_path), exist_ok=True)
                
                # 打开文档并转换
                try:
                    doc = word.Documents.Open(abs_input_path)
                    # PDF格式代码为17
                    doc.SaveAs(abs_output_path, FileFormat=17)
                    doc.Close()
                    word.Quit()
                    
                    if os.path.exists(output_path):
                        logger.info(f"使用Word COM组件完成转换: {output_path}")
                        return output_path
                except Exception as e:
                    logger.warning(f"Word文件处理过程中出错: {str(e)}")
                    if word:
                        try:
                            word.Quit()
                        except:
                            pass
            except Exception as e:
                logger.warning(f"无法使用Word COM组件: {str(e)}")
        
        # 尝试找到LibreOffice路径
        libreoffice_paths = [
            'libreoffice',  # 如果在PATH中
            'soffice',      # 有些系统使用这个命令
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            '/usr/bin/libreoffice',
            '/usr/bin/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice'
        ]
        
        # 尝试使用LibreOffice转换
        for libreoffice_path in libreoffice_paths:
            try:
                logger.info(f"尝试使用LibreOffice路径: {libreoffice_path}")
                subprocess.run([
                    libreoffice_path, '--headless', '--convert-to',
                    'pdf', '--outdir', os.path.dirname(output_path), input_path
                ], check=True, timeout=60)

                # 重命名输出文件
                base_name = os.path.basename(input_path)
                pdf_name = os.path.splitext(base_name)[0] + '.pdf'
                temp_output = os.path.join(os.path.dirname(output_path), pdf_name)

                if os.path.exists(temp_output):
                    if temp_output != output_path:  # 避免不必要的重命名
                        os.rename(temp_output, output_path)
                    logger.info(f"使用LibreOffice完成Word转PDF: {output_path}")
                    return output_path
            except Exception as e:
                logger.warning(f"使用路径 {libreoffice_path} 的LibreOffice转换失败: {str(e)}")
                
        # 尝试使用docx2pdf库
        try:
            # 如果没有安装，先pip install docx2pdf
            from docx2pdf import convert
            
            logger.info("尝试使用docx2pdf库转换")
            convert(input_path, output_path)
            
            if os.path.exists(output_path):
                logger.info(f"使用docx2pdf库完成转换: {output_path}")
                return output_path
        except Exception as e:
            logger.warning(f"docx2pdf转换失败: {str(e)}")
        
        # 后备方案 - 使用简单视觉表示，确保中文正确
        try:
            # 读取文档内容
            doc = Document(input_path)
            content = []

            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)
            
            # 创建临时HTML文件
            temp_html_path = os.path.splitext(output_path)[0] + ".html"
            with open(temp_html_path, 'w', encoding='utf-8') as f:
                # 创建HTML文件，使用多种中文字体引用方法确保兼容性
                html_content = f"""
                <!DOCTYPE html>
                <html>
                <head>
                    <meta charset="UTF-8">
                    <style>
                        @font-face {{
                            font-family: 'NotoSansCJK';
                            src: url('https://cdn.jsdelivr.net/gh/googlefonts/noto-cjk@main/Sans/OTF/Chinese-Simplified/NotoSansCJKsc-Regular.otf') format('opentype');
                            font-weight: normal;
                            font-style: normal;
                        }}
                        @font-face {{
                            font-family: 'NotoSerifCJK'; 
                            src: url('https://cdn.jsdelivr.net/gh/googlefonts/noto-cjk@main/Serif/OTF/Chinese-Simplified/NotoSerifCJKsc-Regular.otf') format('opentype');
                            font-weight: normal;
                            font-style: normal;
                        }}
                        body {{
                            font-family: 'NotoSansCJK', 'NotoSerifCJK', 'Microsoft YaHei', 'SimSun', 'SimHei', Arial, sans-serif;
                            font-size: 12pt;
                            line-height: 1.5;
                            margin: 2cm;
                        }}
                        p {{
                            margin-bottom: 1em;
                            text-align: justify;
                        }}
                        @page {{
                            size: A4;
                            margin: 2cm;
                        }}
                    </style>
                </head>
                <body>
                """
                
                # 添加段落
                for para in content:
                    html_content += f"<p>{para}</p>\n"
                
                html_content += """
                </body>
                </html>
                """
                f.write(html_content)
            
            # 使用Chrome/Edge浏览器直接打印为PDF（最可靠的方法）
            browser_success = False
            if os.name == 'nt':  # Windows
                chrome_paths = [
                    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                    r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
                    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
                    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
                    r"C:\Users\%USERNAME%\AppData\Local\Google\Chrome\Application\chrome.exe",
                    r"%LOCALAPPDATA%\Google\Chrome\Application\chrome.exe",
                    r"%PROGRAMFILES%\Google\Chrome\Application\chrome.exe",
                    r"%PROGRAMFILES(X86)%\Google\Chrome\Application\chrome.exe",
                    r"%PROGRAMFILES%\Microsoft\Edge\Application\msedge.exe",
                    r"%PROGRAMFILES(X86)%\Microsoft\Edge\Application\msedge.exe"
                ]
                
                # 解析可能包含环境变量的路径
                parsed_paths = []
                for path in chrome_paths:
                    try:
                        if '%' in path:
                            expanded_path = os.path.expandvars(path)
                            parsed_paths.append(expanded_path)
                        else:
                            parsed_paths.append(path)
                    except:
                        parsed_paths.append(path)
                
                # 尝试所有可能的浏览器路径
                for chrome_path in parsed_paths:
                    if os.path.exists(chrome_path):
                        try:
                            logger.info(f"使用浏览器 {chrome_path} 打印PDF")
                            abs_html_path = os.path.abspath(temp_html_path)
                            abs_pdf_path = os.path.abspath(output_path)
                            
                            # 尝试执行命令将HTML打印为PDF
                            cmd = [
                                chrome_path,
                                "--headless",
                                "--disable-gpu",
                                "--no-margins",
                                "--print-to-pdf-no-header",
                                "--print-to-pdf-no-footer",
                                f"--print-to-pdf={abs_pdf_path}",
                                "file:///" + abs_html_path.replace('\\', '/')
                            ]
                            
                            process = subprocess.run(cmd, check=True, timeout=60, 
                                                   stderr=subprocess.PIPE, stdout=subprocess.PIPE)
                            
                            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                                logger.info(f"使用浏览器打印PDF成功: {output_path}")
                                os.remove(temp_html_path)
                                browser_success = True
                                return output_path
                            else:
                                logger.warning(f"浏览器生成的PDF文件为空或不存在")
                        except Exception as e:
                            logger.warning(f"使用浏览器 {chrome_path} 打印PDF失败: {str(e)}")
            
            # 如果浏览器方法失败，尝试其他方法
     # 如果HTML临时文件还存在，清理它
            if os.path.exists(temp_html_path):
                try:
                    os.remove(temp_html_path)
                except:
                    pass
            
            # 如果到这里，说明所有HTML方法都失败了，使用fpdf2简单方案
            logger.warning("所有HTML转换方法都失败，使用简单ASCII兼容模式")
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            
            # 检查是否为fpdf2
            is_fpdf2 = hasattr(FPDF, 'add_font_from_name')
            
            if is_fpdf2:
                # 尝试添加中文字体
                font_names = ["Microsoft YaHei", "SimSun", "NSimSun", "SimHei"]
                font_added = False
                
                for font_name in font_names:
                    try:
                        pdf.add_font_from_name(font_name, uni=True)
                        pdf.set_font(font_name, size=12)
                        font_added = True
                        logger.info(f"成功加载中文字体: {font_name}")
                        break
                    except Exception as e:
                        logger.warning(f"无法加载字体 {font_name}: {str(e)}")
                
                if not font_added:
                    pdf.set_font("Arial", size=12)
                    logger.warning("无法加载中文字体，使用系统默认字体")
                
                # 添加文本内容
                for para in content:
                    if para.strip():
                        pdf.multi_cell(0, 10, para)
                        pdf.ln(5)  # 段落间距
            else:
                # 普通模式，非中文兼容
                pdf.set_font("Arial", size=11)
                # 分行处理文本，过滤非ASCII字符
                for para in content:
                    if para.strip():
                        # 过滤非ASCII字符替换为问号
                        filtered_line = ''.join(char if ord(char) < 128 else '?' for char in para)
                        if filtered_line.strip():
                            pdf.multi_cell(0, 8, filtered_line)
                        pdf.ln(4)  # 段落间距
            
            # 保存PDF
            pdf.output(output_path)
            logger.info(f"使用{'fpdf2' if is_fpdf2 else 'fpdf'}完成Word转PDF: {output_path}")
            return output_path
        except Exception as e:
            logger.error(f"简化方法转换失败: {str(e)}")
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                # 如果已生成文件，即使有错误也返回
                logger.info(f"尽管有错误但文件已创建: {output_path}")
                return output_path
            raise
    except Exception as e:
        logger.error(f"Word转PDF失败: {str(e)}")
        raise ValueError(f"Word转PDF失败: {str(e)}")
        
    # 确保处理所有可能的返回路径
    if isinstance(output_path, dict) and "output_path" in output_path:
        return output_path["output_path"]
    return output_path


# Excel 转 PDF
def excel_to_pdf(input_path: str, output_path: str, quality: int) -> str:
    """将Excel转换为PDF格式"""
    try:
        logger.info(f"开始Excel转PDF转换: {input_path}")
        conversion_success = False
        
        # 在Windows上，优先使用Office COM组件
        if os.name == 'nt':
            try:
                import win32com.client
                import pythoncom
                
                logger.info("尝试使用Excel COM组件转换")
                pythoncom.CoInitialize()
                excel = win32com.client.Dispatch("Excel.Application")
                excel.Visible = False
                
                # 使用绝对路径
                abs_input_path = os.path.abspath(input_path)
                abs_output_path = os.path.abspath(output_path)
                
                # 确保输出目录存在
                os.makedirs(os.path.dirname(abs_output_path), exist_ok=True)
                
                # 打开文档并转换
                workbook = excel.Workbooks.Open(abs_input_path)
                workbook.ExportAsFixedFormat(0, abs_output_path)  # 0表示PDF格式
                workbook.Close()
                excel.Quit()
                
                if os.path.exists(output_path):
                    logger.info(f"使用Excel COM组件完成转换: {output_path}")
                    return output_path
            except Exception as e:
                logger.warning(f"无法使用Excel COM组件: {str(e)}")
        
        # 尝试找到LibreOffice路径
        libreoffice_paths = [
            'libreoffice',
            'soffice',
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            '/usr/bin/libreoffice',
            '/usr/bin/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice'
        ]
        
        # 尝试使用LibreOffice转换
        for libreoffice_path in libreoffice_paths:
            try:
                logger.info(f"尝试使用LibreOffice路径: {libreoffice_path}")
                subprocess.run([
                    libreoffice_path, '--headless', '--convert-to',
                    'pdf', '--outdir', os.path.dirname(output_path), input_path
                ], check=True, timeout=60)

                # 重命名输出文件
                base_name = os.path.basename(input_path)
                pdf_name = os.path.splitext(base_name)[0] + '.pdf'
                temp_output = os.path.join(os.path.dirname(output_path), pdf_name)

                if os.path.exists(temp_output):
                    if temp_output != output_path:  # 避免不必要的重命名
                        os.rename(temp_output, output_path)
                    logger.info(f"使用LibreOffice完成Excel转PDF: {output_path}")
                    return output_path
            except Exception as e:
                logger.warning(f"使用路径 {libreoffice_path} 的LibreOffice转换失败: {str(e)}")
        
        # 尝试使用pandas和matplotlib直接渲染表格到PDF
        try:
            logger.info("尝试使用pandas和matplotlib渲染Excel表格")
            import pandas as pd
            import matplotlib.pyplot as plt
            from matplotlib.backends.backend_pdf import PdfPages
            
            # 读取Excel文件
            excel_file = pd.ExcelFile(input_path)
            sheet_names = excel_file.sheet_names
            
            with PdfPages(output_path) as pdf:
                for sheet_name in sheet_names:
                    df = pd.read_excel(input_path, sheet_name=sheet_name)
                    
                    # 设置图表大小
                    if quality == 1:
                        fig_size = (8, 11)  # 低质量，较小的页面
                    elif quality == 2:
                        fig_size = (10, 14)  # 中质量
                    else:
                        fig_size = (12, 16)  # 高质量，较大的页面
                    
                    # 创建图表
                    fig, ax = plt.subplots(figsize=fig_size)
                    ax.axis('tight')
                    ax.axis('off')
                    
                    # 调整表格显示
                    col_width = [0.15] * len(df.columns)
                    if quality > 1:
                        # 优化列宽
                        for i, col in enumerate(df.columns):
                            max_len = max(
                                df[col].astype(str).map(len).max(),
                                len(str(col))
                            ) * 0.01
                            col_width[i] = max(0.1, min(0.3, max_len))
                    
                    # 创建表格
                    table = ax.table(
                        cellText=df.values.tolist(),
                        colLabels=df.columns,
                        loc='center',
                        cellLoc='center',
                        colWidths=col_width
                    )
                    
                    # 设置表格样式
                    table.auto_set_font_size(False)
                    table.set_fontsize(10 if quality > 1 else 8)
                    table.scale(1, 1.5)  # 设置行高
                    
                    # 添加表格标题
                    plt.title(sheet_name)
                    plt.tight_layout()
                    
                    # 保存页面
                    pdf.savefig(fig)
                    plt.close()
            
            if os.path.exists(output_path):
                logger.info(f"使用pandas和matplotlib完成Excel转PDF: {output_path}")
                return output_path
        except Exception as e:
            logger.warning(f"使用pandas和matplotlib转换失败: {str(e)}")
        
        # 如果所有方法都失败，尝试最简单的方法
        logger.warning("所有转换方法失败，尝试简单文本转换")
        try:
            import pandas as pd
            from fpdf import FPDF
            
            # 读取所有工作表
            sheets = pd.read_excel(input_path, sheet_name=None)
            
            # 创建PDF
            pdf = FPDF()
            
            for sheet_name, df in sheets.items():
                pdf.add_page()
                pdf.set_font("Arial", size=12)
                pdf.cell(200, 10, txt=f"Sheet: {sheet_name}", ln=True, align="C")
                
                # 添加列标题
                pdf.set_font("Arial", size=8, style="B")
                col_width = 190 / len(df.columns)  # 自动分配列宽
                for col in df.columns:
                    pdf.cell(col_width, 10, txt=str(col), border=1)
                pdf.ln()
                
                # 添加每行数据
                pdf.set_font("Arial", size=6)
                for _, row in df.iterrows():
                    for val in row:
                        pdf.cell(col_width, 10, txt=str(val), border=1)
                    pdf.ln()
            
            pdf.output(output_path)
            
            if os.path.exists(output_path):
                logger.info(f"使用FPDF完成简单Excel转PDF: {output_path}")
                conversion_success = True
        except Exception as e:
            logger.error(f"简单文本转换失败: {str(e)}")
        
        if not conversion_success:
            raise ValueError("所有Excel转PDF方法均失败")
        
        # 确保处理所有可能的返回路径
        if isinstance(output_path, dict) and "output_path" in output_path:
            return output_path["output_path"]
        return output_path
    except Exception as e:
        logger.error(f"Excel转PDF失败: {str(e)}")
        raise ValueError(f"Excel转PDF失败: {str(e)}")


# PPT 转 PDF
def ppt_to_pdf(input_path: str, output_path: str, quality: int) -> str:
    """将PPT转换为PDF格式"""
    try:
        logger.info(f"开始PPT转PDF转换: {input_path}")
        conversion_success = False
        
        # 在Windows上，优先使用Office COM组件
        if os.name == 'nt':
            try:
                import win32com.client
                import pythoncom
                
                logger.info("尝试使用PowerPoint COM组件转换")
                pythoncom.CoInitialize()
                ppt = win32com.client.Dispatch("PowerPoint.Application")
                ppt.Visible = True
                
                # 使用绝对路径
                abs_input_path = os.path.abspath(input_path)
                abs_output_path = os.path.abspath(output_path)
                
                # 确保输出目录存在
                os.makedirs(os.path.dirname(abs_output_path), exist_ok=True)
                
                # 打开文档并转换
                presentation = ppt.Presentations.Open(abs_input_path, WithWindow=False)
                
                # 固定格式 = 32，PDF = 2
                presentation.ExportAsFixedFormat(abs_output_path, 2, 
                                                PrintRange=0,  # 全部打印
                                                OutputType=0,  # 高质量
                                                PrintHiddenSlides=False,
                                                FrameSlides=True,
                                                Intent=1,      # 屏幕质量
                                                KeepIRMSettings=True)
                presentation.Close()
                ppt.Quit()
                
                if os.path.exists(output_path):
                    logger.info(f"使用PowerPoint COM组件完成转换: {output_path}")
                    return output_path
            except Exception as e:
                logger.warning(f"无法使用PowerPoint COM组件: {str(e)}")
        
        # 尝试找到LibreOffice路径
        libreoffice_paths = [
            'libreoffice',
            'soffice',
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            '/usr/bin/libreoffice',
            '/usr/bin/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice'
        ]
        
        # 尝试使用LibreOffice转换
        for libreoffice_path in libreoffice_paths:
            try:
                logger.info(f"尝试使用LibreOffice路径: {libreoffice_path}")
                subprocess.run([
                    libreoffice_path, '--headless', '--convert-to',
                    'pdf', '--outdir', os.path.dirname(output_path), input_path
                ], check=True, timeout=60)

                # 重命名输出文件
                base_name = os.path.basename(input_path)
                pdf_name = os.path.splitext(base_name)[0] + '.pdf'
                temp_output = os.path.join(os.path.dirname(output_path), pdf_name)

                if os.path.exists(temp_output):
                    if temp_output != output_path:  # 避免不必要的重命名
                        os.rename(temp_output, output_path)
                    logger.info(f"使用LibreOffice完成PPT转PDF: {output_path}")
                    return output_path
            except Exception as e:
                logger.warning(f"使用路径 {libreoffice_path} 的LibreOffice转换失败: {str(e)}")
        
        # 尝试使用python-pptx和reportlab渲染
        try:
            logger.info("尝试使用python-pptx和reportlab渲染PPT")
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            import io
            from reportlab.lib.utils import ImageReader
            from PIL import Image
            
            # 读取PPT
            prs = Presentation(input_path)
            
            # 创建PDF
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter
            
            for slide in prs.slides:
                # 为每个幻灯片创建一个图像
                img_stream = io.BytesIO()
                
                # 使用PIL创建一个白色背景图像
                img = Image.new('RGB', (800, 600), 'white')
                
                # 为每个形状渲染文本
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        # 这里只是简单地将文本添加到图像中
                        from PIL import ImageDraw, ImageFont
                        draw = ImageDraw.Draw(img)
                        try:
                            # 尝试加载一个字体
                            font = ImageFont.truetype("arial.ttf", 14)
                        except:
                            # 如果无法加载，使用默认字体
                            font = ImageFont.load_default()
                        
                        # 绘制文本
                        x, y = 50, 50  # 简单起见，使用固定位置
                        draw.text((x, y), shape.text, fill="black", font=font)
                
                # 保存图像到流
                img.save(img_stream, format='PNG')
                img_stream.seek(0)
                
                # 将图像添加到PDF
                c.drawImage(ImageReader(img_stream), 0, 0, width, height)
                c.showPage()
            
            c.save()
            
            if os.path.exists(output_path):
                logger.info(f"使用python-pptx和reportlab完成PPT转PDF: {output_path}")
                conversion_success = True
        except Exception as e:
            logger.warning(f"使用python-pptx和reportlab转换失败: {str(e)}")
        
        if not conversion_success:
            raise ValueError("所有PPT转PDF方法均失败")
            
        # 确保处理所有可能的返回路径
        if isinstance(output_path, dict) and "output_path" in output_path:
            return output_path["output_path"]
        return output_path
    except Exception as e:
        logger.error(f"PPT转PDF失败: {str(e)}")
        raise ValueError(f"PPT转PDF失败: {str(e)}")


# TXT 转 PDF
def txt_to_pdf(input_path: str, output_path: str, quality: int) -> str:
    """将文本文件转换为PDF格式"""
    try:
        logger.info(f"开始文本转PDF转换: {input_path}")
        
        # 读取文本文件，尝试不同编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030']
        text_content = None
        
        for encoding in encodings:
            try:
                with open(input_path, 'r', encoding=encoding) as file:
                    text_content = file.read()
                logger.info(f"成功使用{encoding}编码读取文本文件")
                break
            except UnicodeDecodeError:
                continue
        
        if text_content is None:
            logger.warning("无法以任何常见编码读取文本文件，尝试二进制读取并解码")
            with open(input_path, 'rb') as file:
                binary_content = file.read()
                # 尝试检测编码
                try:
                    import chardet
                    result = chardet.detect(binary_content)
                    encoding = result['encoding']
                    logger.info(f"检测到编码: {encoding}")
                    text_content = binary_content.decode(encoding)
                except ImportError:
                    logger.warning("未安装chardet，无法自动检测编码")
                    # 尝试使用latin-1，它可以解码任何字节序列
                    text_content = binary_content.decode('latin-1')
        
        # 生成PDF文件
        
        # 方法1：使用reportlab (高质量排版)
        try:
            if quality > 1:  # 中高质量使用reportlab
                logger.info("使用ReportLab生成高质量PDF")
                from reportlab.lib.pagesizes import letter
                from reportlab.lib import colors
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.pdfbase import pdfmetrics
                from reportlab.pdfbase.ttfonts import TTFont
                
                # 尝试注册中文字体
                try:
                    # 尝试加载系统中文字体
                    import platform
                    system = platform.system()
                    
                    if system == "Windows":
                        font_paths = [
                            r"C:\Windows\Fonts\simhei.ttf",    # 黑体
                            r"C:\Windows\Fonts\simsun.ttc",    # 宋体
                            r"C:\Windows\Fonts\msyh.ttc"       # 微软雅黑
                        ]
                    elif system == "Darwin":  # macOS
                        font_paths = [
                            "/Library/Fonts/Arial Unicode.ttf",
                            "/System/Library/Fonts/PingFang.ttc"
                        ]
                    else:  # Linux
                        font_paths = [
                            "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
                            "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc"
                        ]
                    
                    font_registered = False
                    for font_path in font_paths:
                        try:
                            if os.path.exists(font_path):
                                if font_path.endswith('simhei.ttf'):
                                    pdfmetrics.registerFont(TTFont('SimHei', font_path))
                                    logger.info(f"已注册黑体字体")
                                    font_registered = True
                                elif font_path.endswith('simsun.ttc'):
                                    pdfmetrics.registerFont(TTFont('SimSun', font_path))
                                    logger.info(f"已注册宋体字体")
                                    font_registered = True
                                elif font_path.endswith('msyh.ttc'):
                                    pdfmetrics.registerFont(TTFont('MicrosoftYaHei', font_path))
                                    logger.info(f"已注册微软雅黑字体")
                                    font_registered = True
                                elif "noto" in font_path.lower():
                                    pdfmetrics.registerFont(TTFont('NotoSans', font_path))
                                    logger.info(f"已注册Noto Sans字体")
                                    font_registered = True
                                else:
                                    pdfmetrics.registerFont(TTFont('SystemFont', font_path))
                                    logger.info(f"已注册系统字体: {font_path}")
                                    font_registered = True
                                break
                        except Exception as e:
                            logger.warning(f"注册字体失败 {font_path}: {str(e)}")
                    
                    if not font_registered:
                        # 使用默认字体
                        logger.warning("未能注册系统中文字体，将使用默认字体")
                except Exception as e:
                    logger.warning(f"字体注册过程中发生错误: {str(e)}")
                
                # 创建文档
                doc = SimpleDocTemplate(
                    output_path,
                    pagesize=letter,
                    rightMargin=72,
                    leftMargin=72,
                    topMargin=72,
                    bottomMargin=72
                )
                
                # 样式
                styles = getSampleStyleSheet()
                
                # 检查是否已注册中文字体
                font_name = None
                for name in ['SimHei', 'SimSun', 'MicrosoftYaHei', 'NotoSans', 'SystemFont']:
                    if name in pdfmetrics._fonts:
                        font_name = name
                        break
                
                # 创建自定义样式，用于中文文本
                if font_name:
                    chinese_style = ParagraphStyle(
                        'ChineseStyle',
                        parent=styles['Normal'],
                        fontName=font_name,
                        fontSize=12,
                        leading=14,
                        firstLineIndent=20
                    )
                else:
                    chinese_style = styles['Normal']
                
                # 将文本分割成段落
                paragraphs = text_content.split('\n')
                
                # 准备文档内容
                story = []
                for para in paragraphs:
                    if para.strip():  # 如果段落不为空
                        p = Paragraph(para, chinese_style)
                        story.append(p)
                    else:  # 空行用间隔代替
                        story.append(Spacer(1, 12))
                
                # 生成PDF
                doc.build(story)
                logger.info(f"使用ReportLab完成PDF生成: {output_path}")
                return output_path
                
        except Exception as e:
            logger.warning(f"使用ReportLab生成PDF失败: {str(e)}")
        
        # 方法2：使用fpdf (简单但可靠)
        try:
            logger.info("使用FPDF生成PDF")
            from fpdf import FPDF
            
            # 创建PDF实例
            pdf = FPDF()
            pdf.add_page()
            
            # 设置字体 - 尝试支持中文
            try:
                # 尝试使用中文兼容字体
                pdf.add_font('NotoSans', '', r'C:\Windows\Fonts\simhei.ttf', uni=True)
                pdf.set_font('NotoSans', '', 11)
                logger.info("使用中文字体")
            except Exception as e:
                logger.warning(f"添加中文字体失败: {str(e)}")
                # 退回到基本ASCII
                pdf.set_font('Arial', '', 11)
                logger.warning("使用ASCII兼容字体")
                # 过滤非ASCII字符
                text_content = ''.join(char for char in text_content if ord(char) < 128)
                logger.warning("已过滤非ASCII字符")
            
            # 设置页面间距
            pdf.set_margins(20, 20, 20)
            
            # 分割文本为行
            lines = text_content.split('\n')
            
            # 添加文本
            for line in lines:
                pdf.multi_cell(0, 8, line)
            
            # 保存PDF
            pdf.output(output_path)
            logger.info(f"使用FPDF完成PDF生成: {output_path}")
            
            # 检查文件是否生成
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                return output_path
            else:
                logger.error("FPDF生成的文件大小为0或不存在")
                raise ValueError("生成的PDF文件无效")
                
        except Exception as e:
            logger.error(f"使用FPDF生成PDF失败: {str(e)}")
            raise ValueError(f"文本转PDF失败: {str(e)}")
            
        # 确保处理所有可能的返回路径
        if isinstance(output_path, dict) and "output_path" in output_path:
            return output_path["output_path"]
        return output_path
    except Exception as e:
        logger.error(f"文本转PDF失败: {str(e)}")
        raise ValueError(f"文本转PDF失败: {str(e)}")


# Markdown 转 PDF
def markdown_to_pdf(input_path: str, output_path: str, quality: int) -> str:
    """将Markdown转换为PDF格式"""
    try:
        logger.info(f"开始Markdown转PDF转换: {input_path}")
        
        # 读取Markdown文件
        with open(input_path, 'r', encoding='utf-8') as file:
            markdown_text = file.read()
        
        # 转换为HTML
        html_content = markdown.markdown(markdown_text, extensions=['tables', 'fenced_code'])
        
        # 添加基本样式
        styled_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <style>
                body {{
                    font-family: Arial, 'Microsoft YaHei', sans-serif;
                    line-height: 1.6;
                    margin: 40px;
                    max-width: 800px;
                    margin: 0 auto;
                    padding: 20px;
                }}
                h1, h2, h3, h4, h5, h6 {{
                    color: #333;
                    margin-top: 20px;
                }}
                p {{
                    margin-bottom: 16px;
                }}
                code {{
                    background-color: #f5f5f5;
                    border-radius: 3px;
                    padding: 2px 5px;
                }}
                pre {{
                    background-color: #f5f5f5;
                    border-radius: 3px;
                    padding: 16px;
                    overflow: auto;
                }}
                blockquote {{
                    border-left: 5px solid #ddd;
                    padding-left: 15px;
                    color: #555;
                }}
                table {{
                    border-collapse: collapse;
                    width: 100%;
                }}
                table, th, td {{
                    border: 1px solid #ddd;
                }}
                th, td {{
                    padding: 8px;
                    text-align: left;
                }}
                tr:nth-child(even) {{
                    background-color: #f2f2f2;
                }}
            </style>
        </head>
        <body>
            {html_content}
        </body>
        </html>
        """
        
        # 使用weasyprint生成PDF
        html = HTML(string=styled_html)
        html.write_pdf(output_path)
        
        logger.info(f"Markdown转PDF完成: {output_path}")
        
        # 确保处理所有可能的返回路径
        if isinstance(output_path, dict) and "output_path" in output_path:
            return output_path["output_path"]
        return output_path
    except Exception as e:
        logger.error(f"Markdown转PDF失败: {str(e)}")
        raise ValueError(f"Markdown转PDF失败: {str(e)}")


# 简易Markdown解析函数
def simple_md_to_text(md_text):
    """简单解析Markdown文本为普通文本"""
    import re
    # ... [rest of the code remains unchanged]